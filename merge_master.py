import glob
import os
from docx import Document  # For Word
from pypdf import PdfWriter  # For PDF
from pptx import Presentation  # For PowerPoint

# --- Configuration ---

# 1. Set the type of merge you want to perform.
#    Options: 'word', 'pdf', or 'ppt'
MERGE_TYPE = 'ppt'

# 2. This is the folder where ALL your files are located.
SOURCE_FOLDER = "files_to_merge"

# 3. This is the base name for the final merged file.
#    The script will add '.docx', '.pdf', or '.pptx' automatically.
OUTPUT_FILENAME_BASE = "master_document"


# --- End of Configuration ---


def get_merge_files(directory, file_extension):
    """
    Finds all files with a specific extension in the source directory 
    and sorts them alphabetically.
    """
    search_path = os.path.join(directory, f"*.{file_extension}")
    files = sorted(glob.glob(search_path))
    return files


def merge_word_documents(files_list, output_filename):
    """Merges a list of .docx files with page breaks."""
    if not files_list:
        print(f"No .docx files found in '{SOURCE_FOLDER}'.")
        return

    print(f"Found {len(files_list)} Word files to merge.")
    master_doc = Document()
    total_files = len(files_list)

    for i, file_path in enumerate(files_list):
        print(f"Merging: {file_path}")
        source_doc = Document(file_path)
        for element in source_doc.element.body:
            master_doc.element.body.append(element)
        if i < total_files - 1:
            print("  -> Adding page break...")
            master_doc.add_page_break()

    try:
        master_doc.save(output_filename)
        print(f"\nSuccessfully merged {total_files} files into '{output_filename}'")
    except Exception as e:
        print(f"\nAn error occurred while saving the file: {e}")


def merge_pdf_documents(files_list, output_filename):
    """Merges a list of .pdf files."""
    if not files_list:
        print(f"No .pdf files found in '{SOURCE_FOLDER}'.")
        return

    print(f"Found {len(files_list)} PDF files to merge.")
    merger = PdfWriter()
    total_files = len(files_list)

    for i, file_path in enumerate(files_list):
        print(f"Merging: {file_path}")
        try:
            merger.append(file_path)
        except Exception as e:
            print(f"  -> ERROR: Could not merge '{file_path}'. Skipping. Error: {e}")

    try:
        print("\nSaving final document...")
        merger.write(output_filename)
        merger.close()
        print(f"Successfully merged {total_files} files into '{output_filename}'")
    except Exception as e:
        print(f"\nAn error occurred while saving the file: {e}")


def merge_ppt_documents(files_list, output_filename):
    """
    Merges a list of .pptx files.

    WARNING: This is a best-effort merge. It copies slides and text
    from placeholders. It will NOT perfectly preserve:
    - Images
    - Charts, Tables, or other complex objects
    - Custom shapes (e.g., arrows, text boxes not in a placeholder)
    - Formatting (colors, fonts, etc.) if slide layouts differ.
    """
    if not files_list:
        print(f"No .pptx files found in '{SOURCE_FOLDER}'.")
        return

    print(f"Found {len(files_list)} PowerPoint files to merge.")

    # The first presentation's layouts will be the master layouts
    master_prs = Presentation(files_list[0])

    # Create a map of layout names from the master presentation
    layout_name_map = {layout.name: layout for layout in master_prs.slide_layouts}

    # We already "merged" the first file, so start from the second
    for i, file_path in enumerate(files_list[1:]):
        print(f"Merging: {file_path}")
        source_prs = Presentation(file_path)

        for slide in source_prs.slides:
            # Try to find a matching slide layout by name
            source_layout_name = slide.slide_layout.name
            master_layout = layout_name_map.get(source_layout_name)

            if master_layout is None:
                # Fallback to a default layout (e.g., Title and Content)
                print(f"  -> Warning: Layout '{source_layout_name}' not in master. Using default.")
                master_layout = master_prs.slide_layouts[1]  # Index 1 is 'Title and Content'

            # Add the new slide
            new_slide = master_prs.slides.add_slide(master_layout)

            # Copy placeholder text
            for ph in slide.placeholders:
                if not ph.is_placeholder:
                    continue
                try:
                    new_ph = new_slide.placeholders[ph.placeholder_format.idx]
                    if ph.has_text_frame:
                        new_ph.text = ph.text
                except (KeyError, AttributeError):
                    print(f"  -> Warning: Could not copy placeholder index {ph.placeholder_format.idx}")

            # Basic title copy
            if slide.shapes.title:
                if new_slide.shapes.title:
                    new_slide.shapes.title.text = slide.shapes.title.text

    try:
        master_prs.save(output_filename)
        print(f"\nSuccessfully merged {len(files_list)} files into '{output_filename}'")
    except Exception as e:
        print(f"\nAn error occurred while saving the file: {e}")


# --- Main execution ---
if __name__ == "__main__":
    if not os.path.exists(SOURCE_FOLDER):
        os.makedirs(SOURCE_FOLDER)
        print(f"Created source directory: '{SOURCE_FOLDER}'")
        print("Please add your files to this folder and run the script again.")
    else:

        if MERGE_TYPE.lower() == 'word':
            print("--- Starting WORD merge ---")
            files = get_merge_files(SOURCE_FOLDER, "docx")
            output_file = f"{OUTPUT_FILENAME_BASE}.docx"
            merge_word_documents(files, output_file)

        elif MERGE_TYPE.lower() == 'pdf':
            print("--- Starting PDF merge ---")
            files = get_merge_files(SOURCE_FOLDER, "pdf")
            output_file = f"{OUTPUT_FILENAME_BASE}.pdf"
            merge_pdf_documents(files, output_file)

        elif MERGE_TYPE.lower() == 'ppt':
            print("--- Starting POWERPOINT (.pptx) merge ---")
            files = get_merge_files(SOURCE_FOLDER, "pptx")
            output_file = f"{OUTPUT_FILENAME_BASE}.pptx"
            merge_ppt_documents(files, output_file)

        else:
            print(f"Error: Invalid MERGE_TYPE: '{MERGE_TYPE}'")
            print("Please open the script and set MERGE_TYPE to 'word', 'pdf', or 'ppt'.")
